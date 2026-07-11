from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    # Initialize presentation
    prs = Presentation()
    
    # Set to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use a blank slide layout
    blank_layout = prs.slide_layouts[6]

    # --- Utility Functions ---
    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        return title_box

    def add_bullet_points(slide, bullets, left, top, width, height):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for i, text in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            # Split bold headers if applicable
            if ":" in text and len(text.split(":")[0]) < 40:
                header, body = text.split(":", 1)
                run1 = p.add_run()
                run1.text = header + ":"
                run1.font.bold = True
                run1.font.size = Pt(20)
                run2 = p.add_run()
                run2.text = body
                run2.font.size = Pt(20)
            else:
                p.text = text
                p.font.size = Pt(20)
            p.level = 0
            p.space_after = Pt(14)
            p.line_spacing = 1.1

    def add_table(slide, rows, cols, left, top, width, height, data):
        shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = shape.table
        
        for r_idx, row_data in enumerate(data):
            for c_idx, cell_data in enumerate(row_data):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(cell_data)
                
                # Format cell text
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(16)
                        if r_idx == 0:
                            run.font.bold = True
        return table

    # ==========================================
    # Slide 1: ML Baseline
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    add_title(slide1, "ML Baseline: Pure Causal Replication & Sandbox Limits")
    
    s1_bullets = [
        "Pristine Replication: Rebuilt 100% vanilla Inverse Probability Weighting (IPW) framework to match JAE paper benchmarks exactly.",
        "Zero ML Heuristics: Eliminated class_weight='balanced' and SMOTE to guarantee pure econometric uncalibrated consistency.",
        "The Overfitting Ceiling: Proved advanced ML hacks (Elastic Net, OOF Meta-Fusion) collapse dense embedding manifolds in small-N, rare-event regimes (N ≈ 400, base rate ≈ 5%)."
    ]
    add_bullet_points(slide1, s1_bullets, Inches(0.5), Inches(1.5), Inches(6), Inches(5))
    
    s1_table_data = [
        ["Target", "Metric", "Model", "Paper"],
        ["Unweighted", "NDCG@K", "57.85%", "57.27%"],
        ["Weighted", "NDCG@K", "64.48%", "64.48%"],
        ["Both", "Precision@K", "50.00%", "50.00%"]
    ]
    add_table(slide1, 4, 4, Inches(6.8), Inches(2.5), Inches(6), Inches(2), s1_table_data)

    # ==========================================
    # Slide 2: Entity Alignment
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, "Cold-Start Entity Alignment & Temporal Protection")
    
    s2_bullets = [
        "The Alignment Trap: Filtered CSV row indices (439-569) were fully anonymized; raw PDF directories used disjoint, chaotic file-naming.",
        "Alphabetical Zip Discovery: Uncovered that sequential alphabetical sorting of raw PDF folders perfectly mirrors the 1:1 progression of the CSV row progression.",
        "Deterministic Mapping: Locked down a flawless candidate_mapping.csv matrix to unblock both prompting and peer fine-tuning tracks.",
        "Temporal Protection Layer: Parsed JMP text strictly up to jmp_intro_char_limit (first 10 pages), physically stripping Appendices and References to eliminate Look-Ahead Bias."
    ]
    add_bullet_points(slide2, s2_bullets, Inches(0.5), Inches(1.5), Inches(12), Inches(5))

    # ==========================================
    # Slide 3: LLM Benchmarks
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, "Dual-Track Zero-Shot LLM Benchmarks & The 'Semantic Trap'")
    
    s3_bullets = [
        "Track A (GPT-2 Large): Collapsed into random noise (AUC ≈ 49.60%) due to strict 1,024-token context constraints truncating the JMP.",
        "Track B (Gemini Flash): Successfully leveraged anti-bias CoT and name scrubbing to extract granular signals, but hit a strict 16.67% Precision ceiling.",
        "The Semantic Trap: Zero-shot LLMs over-index on writing rhetoric and topic interest. Long-term productivity is driven by latent causal graphs (Placerank, co-authorship topologies) captured only by our IPW baseline."
    ]
    add_bullet_points(slide3, s3_bullets, Inches(0.5), Inches(1.3), Inches(12), Inches(3.5))
    
    s3_table_data = [
        ["Metric", "Vanilla IPW", "Track A (GPT-2)", "Track B (Gemini)"],
        ["Precision@K", "50.00%", "16.67%", "16.67%"],
        ["NDCG@K", "64.48%", "15.13%", "15.13%"],
        ["LIFT", "1091.67%", "363.89%", "363.89%"],
        ["AUC", "81.33%", "49.60%", "56.73%"],
        ["F1-Score", "50.00%", "0.00%", "16.67%"]
    ]
    add_table(slide3, 6, 4, Inches(1.5), Inches(4.5), Inches(10), Inches(2.5), s3_table_data)

    # ==========================================
    # Slide 4: Next Steps
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_title(slide4, "Next Steps: System Optimization & Architectural Scaling")
    
    s4_bullets = [
        "Few-Shot Prompting Transition: Inject historical 2017 text anchors (1 True Positive, 1 True Negative) to teach Gemini cross-candidate calibration.",
        "Unblocking Peer Finetuning: Provide the master candidate_mapping.csv to map the raw text data directly to the outcome classification heads.",
        "Hybrid Fusion Modeling: Extract Gemini's continuous text confidence logits and feed them back into the IPW model as a structured covariate to test if Semantic Intuition + Causal Weighting can shatter the 64.48% baseline."
    ]
    add_bullet_points(slide4, s4_bullets, Inches(0.5), Inches(1.5), Inches(12), Inches(5))

    # Save presentation
    output_filename = "Week4_Progress_Update_Slides.pptx"
    prs.save(output_filename)
    print(f"Presentation generated successfully: {os.path.abspath(output_filename)}")

if __name__ == "__main__":
    create_presentation()
