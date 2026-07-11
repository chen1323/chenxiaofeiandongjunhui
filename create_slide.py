from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_capstone_slide():
    prs = Presentation()
    # 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Blank slide layout
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Colors
    navy_blue = RGBColor(31, 78, 121) # 1F4E79
    mid_blue = RGBColor(46, 117, 182) # 2E75B6
    dark_gray = RGBColor(64, 64, 64)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Week 5-6 Progress Update: Hybrid Architecture Grid Sweep & System Blueprint"
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = navy_blue
    
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = mid_blue
    line.line.color.rgb = mid_blue
    
    # --- Column 1: Experiments ---
    txBox1 = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.0), Inches(5.5))
    tf1 = txBox1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "1. Experiments & The \"Semantic Trap\" Verification"
    p1.font.name = 'Arial'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = navy_blue
    
    bullet_font_size = Pt(16)
    
    bullets_col1 = [
        ("Data Integrity:", " Secured the 2018 mapping matrix via Alphabetical Zip Logic, ensuring a zero-leakage test environment."),
        ("Hybrid Architectures:", " Engineered LLM-driven 'Advanced Referee Prompting' (Originality, Methodology, Empirical Rigor, Relevance) and injected logit rewards (+1.0 to +2.0) into the econometric baseline."),
        ("Key Finding (The Semantic Trap):", " Escalating semantic rewards caused a catastrophic NDCG@K collapse (from 40.34% to 33.90%) at the extreme upper tail (K=8)."),
        ("Conclusion:", " Pure textual rigor cannot override the structural causal network (Placerank, co-author momentum). The Vanilla IPW model remains the undisputed global maximum.")
    ]
    
    for bold_text, normal_text in bullets_col1:
        p = tf1.add_paragraph()
        p.level = 0
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = bold_text
        run1.font.name = 'Arial'
        run1.font.size = bullet_font_size
        run1.font.bold = True
        run1.font.color.rgb = dark_gray
        
        run2 = p.add_run()
        run2.text = normal_text
        run2.font.name = 'Arial'
        run2.font.size = bullet_font_size
        run2.font.color.rgb = dark_gray
        
    # --- Column 2: System Architecture & Next Steps ---
    txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6.0), Inches(5.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    p2 = tf2.paragraphs[0]
    p2.text = "2. Proposed System Architecture (For Discussion)"
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = navy_blue
    
    bullets_col2_1 = [
        ("Current Status:", " Finalizing model lock; pivoting to system deployment planning."),
        ("System Blueprint:", " Designing a Stateful Multi-Step UI flowchart to illustrate the complete operational pipeline."),
        ("Dual-Track Routing Concept:", "")
    ]
    
    for bold_text, normal_text in bullets_col2_1:
        p = tf2.add_paragraph()
        p.level = 0
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = bold_text
        run1.font.name = 'Arial'
        run1.font.size = bullet_font_size
        run1.font.bold = True
        run1.font.color.rgb = dark_gray
        
        run2 = p.add_run()
        run2.text = normal_text
        run2.font.name = 'Arial'
        run2.font.size = bullet_font_size
        run2.font.color.rgb = dark_gray
        
    # Sub-bullets for Dual-Track
    sub_bullets = [
        ("Enterprise Track: Batch evaluation with Semantic Tag Diminution.", 1),
        ("Applicant Track: Explainable AI (XAI) dashboard with counterfactual feedback loops.", 1)
    ]
    for text, level in sub_bullets:
        p = tf2.add_paragraph()
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.name = 'Arial'
        run.font.size = bullet_font_size
        run.font.color.rgb = dark_gray

    # Section 3
    p3 = tf2.add_paragraph()
    p3.text = "3. Next Steps"
    p3.space_before = Pt(24)
    p3.font.name = 'Arial'
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = navy_blue
    
    bullets_col2_2 = [
        ("Model:", " Lock the analytical repository and Vanilla IPW framework."),
        ("Drafting:", " Translate the 'Semantic Trap' findings into the core Discussion chapter of the thesis."),
        ("Engineering:", " Finalize the full system flowchart design to guide the upcoming web UI construction.")
    ]
    for bold_text, normal_text in bullets_col2_2:
        p = tf2.add_paragraph()
        p.level = 0
        p.space_before = Pt(10)
        run1 = p.add_run()
        run1.text = bold_text
        run1.font.name = 'Arial'
        run1.font.size = bullet_font_size
        run1.font.bold = True
        run1.font.color.rgb = dark_gray
        
        run2 = p.add_run()
        run2.text = normal_text
        run2.font.name = 'Arial'
        run2.font.size = bullet_font_size
        run2.font.color.rgb = dark_gray
        
    prs.save("Week5_6_Capstone_Update.pptx")

if __name__ == "__main__":
    create_capstone_slide()
