import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
import re

# 1. DESIGN CONSTANTS (Based on academic-pptx-skill)
COLORS = {
    "bg": RGBColor(0xFF, 0xFF, 0xFF),
    "primary": RGBColor(0x1F, 0x4E, 0x79),   # Dark navy — titles
    "accent": RGBColor(0x2E, 0x75, 0xB6),    # Mid-blue — headers, highlights
    "body": RGBColor(0x2D, 0x2D, 0x2D),      # Near-black — body text
    "rule": RGBColor(0xCC, 0xCC, 0xCC),      # Light gray — divider lines
    "muted": RGBColor(0x77, 0x77, 0x77),     # Gray — citations, captions
}

FONTS = {
    "face": "Arial",
    "title": 26,
    "sectionHeader": 20,
    "body": 18,
}

def parse_markdown(filepath):
    """
    Parses the Week2_Progress_Update_Slide.md file.
    Expects sections starting with ### and bullets starting with *
    """
    with open(filepath, 'r') as f:
        lines = f.readlines()
    
    sections = []
    current_header = None
    current_bullets = []
    
    for line in lines:
        line = line.strip()
        if line.startswith("###"):
            if current_header:
                sections.append((current_header, current_bullets))
            current_header = line.replace("###", "").replace("**", "").strip()
            current_bullets = []
        elif line.startswith("*"):
            bullet_text = line.replace("*", "").replace("**", "").strip()
            current_bullets.append(bullet_text)
    
    if current_header:
        sections.append((current_header, current_bullets))
        
    return sections

def create_pptx(sections, output_path):
    prs = Presentation()
    # Set to 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["bg"]

    # Action Title (Synthesized Takeaway)
    action_title = "Weeks 1-2: Replicated SMOTE Baseline and Engineered De-biased LLM Pipeline"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.33), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = action_title
    p.font.size = Pt(FONTS["title"])
    p.font.name = FONTS["face"]
    p.font.bold = True
    p.font.color.rgb = COLORS["primary"]
    p.alignment = PP_ALIGN.LEFT

    # Divider Line
    line_shape = slide.shapes.add_shape(
        1, # Rectangle
        Inches(0.5), Inches(1.2), Inches(12.33), Inches(0.02)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = COLORS["rule"]
    line_shape.line.color.rgb = COLORS["rule"]

    # Content Columns or Vertical Stack
    # Since we have 3 sections, we'll stack them vertically with clear spacing
    y_pos = Inches(1.5)
    
    for header, bullets in sections:
        # Section Header
        h_box = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(12.33), Inches(0.4))
        hp = h_box.text_frame.paragraphs[0]
        hp.text = header
        hp.font.size = Pt(FONTS["sectionHeader"])
        hp.font.name = FONTS["face"]
        hp.font.bold = True
        hp.font.color.rgb = COLORS["accent"]
        
        y_pos += Inches(0.5)
        
        # Bullets
        b_box = slide.shapes.add_textbox(Inches(0.7), y_pos, Inches(12.13), Inches(1.5))
        btf = b_box.text_frame
        btf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                bp = btf.paragraphs[0]
            else:
                bp = btf.add_paragraph()
            bp.text = bullet
            bp.font.size = Pt(FONTS["body"])
            bp.font.name = FONTS["face"]
            bp.font.color.rgb = COLORS["body"]
            bp.level = 0
            bp.space_after = Pt(8)
            bp.alignment = PP_ALIGN.LEFT
            
        y_pos += Inches(0.4 * len(bullets)) + Inches(0.2)

    # Footer / Institutional Label
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4))
    fp = footer_box.text_frame.paragraphs[0]
    fp.text = "Capstone Project: Identifying Research Talent using ML/LLMs  |  Status Update: May 2026"
    fp.font.size = Pt(12)
    fp.font.name = FONTS["face"]
    fp.font.color.rgb = COLORS["muted"]

    prs.save(output_path)
    print(f"Generated {output_path}")

if __name__ == "__main__":
    md_file = "Week2_Progress_Update_Slide.md"
    pptx_file = "Week2_Progress_Update.pptx"
    
    if os.path.exists(md_file):
        content_sections = parse_markdown(md_file)
        create_pptx(content_sections, pptx_file)
    else:
        print(f"Error: {md_file} not found.")
